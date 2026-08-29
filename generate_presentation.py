from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen standard dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ==========================================
    # SWISS EDITORIAL PALETTE WITH DISTINCT ACCENT DIFFERENTIATION
    # ==========================================
    BG_WHITE = RGBColor(255, 255, 255)       # Pure White (#FFFFFF)
    
    TEXT_BLACK = RGBColor(17, 17, 17)        # Deep Charcoal Black (#111111)
    TEXT_HEAVY = RGBColor(38, 38, 38)        # Charcoal-800 (#262626)
    TEXT_MUTED = RGBColor(115, 115, 115)     # Editorial Gray (#737373)
    TEXT_SUBTLE = RGBColor(163, 163, 163)    # Hairline Gray (#A3A3A3)

    LINE_COLOR = RGBColor(229, 229, 229)     # Sharp Grid Hairline (#E5E5E5)
    LINE_DARK = RGBColor(17, 17, 17)         # Black Accent Rule (#111111)

    # Distinct Semantic & Architectural Accent Colors
    COLOR_ROSE = RGBColor(225, 29, 72)       # Rose-600 (#E11D48) - Bottlenecks/Alerts
    COLOR_VIOLET = RGBColor(124, 58, 237)    # Violet-600 (#7C3AED) - Backend & APIs
    COLOR_TEAL = RGBColor(13, 148, 136)      # Teal-600 (#0D9488) - Vision & AI
    COLOR_EMERALD = RGBColor(5, 150, 105)    # Emerald-600 (#059669) - Value & Outputs
    COLOR_AMBER = RGBColor(217, 119, 6)      # Amber-600 (#D97706) - Storage & Client
    COLOR_CYAN = RGBColor(2, 132, 199)       # Sky/Cyan-600 (#0284C7) - Network & Traffic
    COLOR_INDIGO = RGBColor(79, 70, 229)     # Indigo-600 (#4F46E5) - Compute & Infrastructure

    def apply_slide_base(slide, slide_num, total_slides=6):
        # 1. Add Smooth Native Slide Transition Animation
        try:
            tr_xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>'
            slide._element.append(parse_xml(tr_xml))
        except Exception:
            pass

        # 2. Pure White Canvas
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_WHITE
        bg.line.color.rgb = BG_WHITE

        # 3. Architectural Top Line
        top_rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.45), Inches(11.733), Inches(0.015))
        top_rule.fill.solid()
        top_rule.fill.fore_color.rgb = LINE_DARK
        top_rule.line.color.rgb = LINE_DARK

        # 4. Architectural Footer Grid
        bot_rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.9), Inches(11.733), Inches(0.01))
        bot_rule.fill.solid()
        bot_rule.fill.fore_color.rgb = LINE_COLOR
        bot_rule.line.color.rgb = LINE_COLOR

        footer_tb = slide.shapes.add_textbox(Inches(0.8), Inches(6.98), Inches(11.733), Inches(0.3))
        ftf = footer_tb.text_frame
        ftf.word_wrap = False
        ftf.margin_left = ftf.margin_top = ftf.margin_right = ftf.margin_bottom = 0
        fp = ftf.paragraphs[0]
        fp.text = f"INTELLIGENT DOCUMENT PROCESSING  //  SYSTEM ARCHITECTURE  //  0{slide_num} OF 0{total_slides}"
        fp.font.size = Pt(8.5)
        fp.font.bold = True
        fp.font.color.rgb = TEXT_MUTED

    def add_hairline(slide, left, top, width, height, color=LINE_COLOR):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.line.color.rgb = color
        return line

    def add_editorial_header(slide, section_tag, title_text, subtitle_text, tag_color=TEXT_MUTED):
        # Section Micro-Label with distinct color
        tag_tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.56), Inches(11.733), Inches(0.25))
        ttf = tag_tb.text_frame
        ttf.word_wrap = False
        ttf.margin_left = ttf.margin_top = ttf.margin_right = ttf.margin_bottom = 0
        tp = ttf.paragraphs[0]
        tp.text = section_tag.upper()
        tp.font.size = Pt(9)
        tp.font.bold = True
        tp.font.color.rgb = tag_color

        # Main Headline
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.86), Inches(11.733), Inches(0.55))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_BLACK

        # Subtitle
        if subtitle_text:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.42), Inches(11.733), Inches(0.35))
            tf_sub = sub_box.text_frame
            tf_sub.word_wrap = True
            tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle_text
            p_sub.font.size = Pt(11)
            p_sub.font.color.rgb = TEXT_MUTED

        # Horizontal Divider below header
        add_hairline(slide, Inches(0.8), Inches(1.82), Inches(11.733), Inches(0.01))

    # ==========================================
    # SLIDE 1: EDITORIAL TYPOGRAPHIC COVER
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide1, 1, 6)

    # Micro Tag
    tb_cov_tag = slide1.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.733), Inches(0.3))
    tf_ct = tb_cov_tag.text_frame
    tf_ct.margin_left = tf_ct.margin_top = tf_ct.margin_right = tf_ct.margin_bottom = 0
    p_ct = tf_ct.paragraphs[0]
    p_ct.text = "✦  SYSTEM SPECIFICATION & ARCHITECTURE  //  2026"
    p_ct.font.size = Pt(10)
    p_ct.font.bold = True
    p_ct.font.color.rgb = COLOR_VIOLET

    # Giant Main Title
    tb_title = slide1.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(11.733), Inches(1.8))
    tf_t = tb_title.text_frame
    tf_t.word_wrap = True
    tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0

    p_t1 = tf_t.paragraphs[0]
    p_t1.text = "Intelligent Document\nProcessing (IDP)"
    p_t1.font.size = Pt(44)
    p_t1.font.bold = True
    p_t1.font.color.rgb = TEXT_BLACK

    # Architectural Mid Divider Line
    add_hairline(slide1, Inches(0.8), Inches(3.65), Inches(11.733), Inches(0.015), color=LINE_DARK)

    # Two Column Layout below title
    tb_lead = slide1.shapes.add_textbox(Inches(0.8), Inches(3.95), Inches(5.6), Inches(2.3))
    tf_lead = tb_lead.text_frame
    tf_lead.word_wrap = True
    tf_lead.margin_left = tf_lead.margin_top = tf_lead.margin_right = tf_lead.margin_bottom = 0

    p_l1 = tf_lead.paragraphs[0]
    p_l1.text = "Transforming Complex Technical & Engineering Documents into Actionable Structured Intelligence"
    p_l1.font.size = Pt(15)
    p_l1.font.bold = True
    p_l1.font.color.rgb = TEXT_HEAVY

    p_l2 = tf_lead.add_paragraph()
    p_l2.space_before = Pt(10)
    p_l2.text = "An enterprise AI platform engineered to eliminate manual transcription bottlenecks from massive O&M manuals, vendor spare parts catalogs (RSPL), engineering schematics, and handwritten shift logs."
    p_l2.font.size = Pt(11.5)
    p_l2.font.color.rgb = TEXT_MUTED

    # Vertical Hairline Divider
    add_hairline(slide1, Inches(6.8), Inches(3.95), Inches(0.01), Inches(2.4))

    # Right Column: 3 Architectural Pillars with Symbols & Color Coding
    tb_caps = slide1.shapes.add_textbox(Inches(7.2), Inches(3.95), Inches(5.3), Inches(2.4))
    tf_caps = tb_caps.text_frame
    tf_caps.word_wrap = True
    tf_caps.margin_left = tf_caps.margin_top = tf_caps.margin_right = tf_caps.margin_bottom = 0

    caps_data = [
        ("⚡ 01 // DUAL-STREAM INGESTION", "Sub-second native PDF text stream & adaptive 300-DPI multimodal OCR vision.", COLOR_VIOLET),
        ("⚙ 02 // DOMAIN-TARGETED SCHEMAS", "Direct extraction of maintenance schedules, spare parts catalogs, and troubleshooting trees.", COLOR_TEAL),
        ("☁ 03 // AWS CLOUD ARCHITECTURE", "Serverless auto-scaling ECS Fargate containers, KMS-encrypted S3, and audit grounding.", COLOR_INDIGO)
    ]

    for idx, (cap_title, cap_desc, cap_col) in enumerate(caps_data):
        p_c_title = tf_caps.paragraphs[0] if idx == 0 else tf_caps.add_paragraph()
        if idx > 0:
            p_c_title.space_before = Pt(10)
        p_c_title.text = cap_title
        p_c_title.font.size = Pt(10.5)
        p_c_title.font.bold = True
        p_c_title.font.color.rgb = cap_col

        p_c_desc = tf_caps.add_paragraph()
        p_c_desc.space_before = Pt(2)
        p_c_desc.text = cap_desc
        p_c_desc.font.size = Pt(10)
        p_c_desc.font.color.rgb = TEXT_MUTED

    slide1.notes_slide.notes_text_frame.text = (
        "Hello everyone. Welcome to this presentation on Intelligent Document Processing (IDP). "
        "Today, I will walk you through how we addressed manual bottlenecks in technical document extraction by replacing "
        "rigid heuristics with a multimodal AI layer, and how this architecture is built and deployed on AWS."
    )

    # ==========================================
    # SLIDE 2: USE CASE & WHAT WE ADDRESS
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide2, 2, 6)
    add_editorial_header(
        slide2,
        "01 // OPERATIONAL CONTEXT & PROBLEM SPACE",
        "From Complex Unstructured Manuals to Actionable Data",
        "Solving critical operational bottlenecks in industrial maintenance and technical documentation",
        tag_color=COLOR_ROSE
    )

    col_w = Inches(3.644)
    gap = Inches(0.4)
    start_x = Inches(0.8)
    grid_y = Inches(2.05)
    grid_h = Inches(4.55)

    cards_data = [
        {
            "num": "01",
            "tag": "THE BOTTLENECK",
            "color": COLOR_ROSE,
            "title": "Manual Transcription",
            "bullets": [
                ("⚠️ Massive Manuals", "200–500+ page PDFs required days of manual review."),
                ("⚠️ High Error Rates", "Typos in complex part codes and inspection intervals."),
                ("⚠️ Format Silos", "Varied, non-standard layouts across equipment vendors."),
                ("⚠️ ERP Disconnect", "Critical maintenance knowledge trapped in static PDF files.")
            ]
        },
        {
            "num": "02",
            "tag": "THE SOLUTION",
            "color": COLOR_TEAL,
            "title": "Targeted AI Extraction",
            "bullets": [
                ("⚙ Maintenance Tasks", "Routines (Daily/Monthly/Annual) & instruction steps."),
                ("⚙ Spare Parts (RSPL)", "Item names, numbers, drawing codes & quantities."),
                ("⚙ Troubleshooting", "Identifies symptoms, probable causes & verified fixes."),
                ("⚙ Schema Normalization", "Auto-mapped to unified equipment taxonomies.")
            ]
        },
        {
            "num": "03",
            "tag": "BUSINESS VALUE",
            "color": COLOR_EMERALD,
            "title": "Operational Impact",
            "bullets": [
                ("✓ 90%+ Turnaround Boost", "Near-instant transformation from PDF to structured data."),
                ("✓ Page-Level Grounding", "Every extracted row is mapped to its exact source page."),
                ("✓ Direct Sync & Export", "One-click export to multi-tab Excel workbooks & databases."),
                ("✓ Human-in-the-Loop", "Integrated editor and verification workflow for QA.")
            ]
        }
    ]

    for i, cdata in enumerate(cards_data):
        cx = start_x + i * (col_w + gap)

        add_hairline(slide2, cx, grid_y, Inches(1.0), Inches(0.04), color=cdata["color"])

        if i < len(cards_data) - 1:
            add_hairline(slide2, cx + col_w + Inches(0.2), grid_y, Inches(0.01), grid_h)

        tb_c = slide2.shapes.add_textbox(cx, grid_y + Inches(0.12), col_w, grid_h - Inches(0.12))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0

        p0 = tf_c.paragraphs[0]
        p0.text = f"{cdata['num']} // {cdata['tag']}"
        p0.font.size = Pt(9.5)
        p0.font.bold = True
        p0.font.color.rgb = cdata["color"]

        p1 = tf_c.add_paragraph()
        p1.space_before = Pt(4)
        p1.text = cdata["title"]
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_BLACK

        for btitle, bdesc in cdata["bullets"]:
            pt = tf_c.add_paragraph()
            pt.space_before = Pt(10)
            pt.text = btitle
            pt.font.size = Pt(11)
            pt.font.bold = True
            pt.font.color.rgb = TEXT_BLACK

            pd = tf_c.add_paragraph()
            pd.space_before = Pt(1)
            pd.text = f"    {bdesc}"
            pd.font.size = Pt(9.5)
            pd.font.color.rgb = TEXT_MUTED

    slide2.notes_slide.notes_text_frame.text = (
        "In our operational journey, technical manuals were the primary bottleneck. Engineers spent countless hours "
        "manually transcribing complex spare parts tables and maintenance checklists. We addressed this directly by building "
        "a domain-aware extraction engine that structures maintenance routines, RSPL parts, and troubleshooting trees with "
        "exact page-level grounding and immediate spreadsheet export."
    )

    # ==========================================
    # SLIDE 3: WORKING OF APP (HEURISTICS -> AI LAYER)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide3, 3, 6)
    add_editorial_header(
        slide3,
        "02 // SYSTEM ARCHITECTURE & DUAL-STREAM INGESTION",
        "How the App Works: From Heuristics to an AI-First Layer",
        "Adaptive dual-stream processing for both clean digital PDFs and degraded/handwritten scans",
        tag_color=COLOR_INDIGO
    )

    # Top Evolution Statement Strip
    tb_evo = slide3.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.733), Inches(0.85))
    tf_evo = tb_evo.text_frame
    tf_evo.word_wrap = True
    tf_evo.margin_left = tf_evo.margin_top = tf_evo.margin_right = tf_evo.margin_bottom = 0

    ep0 = tf_evo.paragraphs[0]
    ep0.text = "PARADIGM EVOLUTION // MOVING BEYOND FRAGILE HEURISTICS"
    ep0.font.size = Pt(9.5)
    ep0.font.bold = True
    ep0.font.color.rgb = COLOR_AMBER

    ep1 = tf_evo.add_paragraph()
    ep1.space_before = Pt(4)
    ep1.text = "⚠️ Previously: Brittle regex and fixed coordinate parsers broke whenever vendor table formats changed.\n✓ Now: Multimodal AI understands layout semantics, reads rotated text, and reliably reconstructs complex tables."
    ep1.font.size = Pt(11)
    ep1.font.bold = True
    ep1.font.color.rgb = TEXT_BLACK

    # Mid Horizontal Divider
    add_hairline(slide3, Inches(0.8), Inches(3.0), Inches(11.733), Inches(0.01))

    # Two Wide Columns
    half_w = Inches(5.666)
    stream_y = Inches(3.18)
    stream_h = Inches(3.5)

    add_hairline(slide3, Inches(6.666), stream_y, Inches(0.01), stream_h)

    # Stream 1: Digital Native
    add_hairline(slide3, Inches(0.8), stream_y, Inches(1.2), Inches(0.04), color=COLOR_VIOLET)

    tb_s1 = slide3.shapes.add_textbox(Inches(0.8), stream_y + Inches(0.1), half_w, stream_h - Inches(0.1))
    tf_s1 = tb_s1.text_frame
    tf_s1.word_wrap = True
    tf_s1.margin_left = tf_s1.margin_top = tf_s1.margin_right = tf_s1.margin_bottom = 0

    sp0 = tf_s1.paragraphs[0]
    sp0.text = "STREAM 01 // NATIVE DIGITAL PDFs"
    sp0.font.size = Pt(9.5)
    sp0.font.bold = True
    sp0.font.color.rgb = COLOR_VIOLET

    sp1 = tf_s1.add_paragraph()
    sp1.space_before = Pt(4)
    sp1.text = "High-Speed Digital Text Stream"
    sp1.font.size = Pt(16)
    sp1.font.bold = True
    sp1.font.color.rgb = TEXT_BLACK

    s1_items = [
        ("⚡ Direct Text Parsing", "Extracts digital font glyphs and embedded tables directly via PyMuPDF/pypdf."),
        ("⚡ Zero OCR Overhead", "Eliminates rasterization latency for blazing-fast token processing."),
        ("⚡ Batch Chunking", "Intelligently handles 500+ page manuals with rolling context windowing."),
        ("⚡ LLM Structuring", "Transforms raw text streams into strict O&M JSON schemas.")
    ]
    for title, desc in s1_items:
        pt = tf_s1.add_paragraph()
        pt.space_before = Pt(8)
        pt.text = title
        pt.font.size = Pt(10.5)
        pt.font.bold = True
        pt.font.color.rgb = TEXT_BLACK

        pd = tf_s1.add_paragraph()
        pd.space_before = Pt(1)
        pd.text = f"    {desc}"
        pd.font.size = Pt(9.5)
        pd.font.color.rgb = TEXT_MUTED

    # Stream 2: Scanned & Handwritten
    add_hairline(slide3, Inches(7.066), stream_y, Inches(1.2), Inches(0.04), color=COLOR_TEAL)

    tb_s2 = slide3.shapes.add_textbox(Inches(7.066), stream_y + Inches(0.1), half_w, stream_h - Inches(0.1))
    tf_s2 = tb_s2.text_frame
    tf_s2.word_wrap = True
    tf_s2.margin_left = tf_s2.margin_top = tf_s2.margin_right = tf_s2.margin_bottom = 0

    s2_p0 = tf_s2.paragraphs[0]
    s2_p0.text = "STREAM 02 // SCANNED & HANDWRITTEN"
    s2_p0.font.size = Pt(9.5)
    s2_p0.font.bold = True
    s2_p0.font.color.rgb = COLOR_TEAL

    s2_p1 = tf_s2.add_paragraph()
    s2_p1.space_before = Pt(4)
    s2_p1.text = "Multimodal Vision & Adaptive OCR"
    s2_p1.font.size = Pt(16)
    s2_p1.font.bold = True
    s2_p1.font.color.rgb = TEXT_BLACK

    s2_items = [
        ("👁️ High-DPI Rasterization", "Renders crystal-clear page tiles with auto-rotation angle detection."),
        ("👁️ Vision AI Reasoning", "Multimodal LLMs decipher handwriting, stamps, and complex catalog drawings."),
        ("👁️ Adaptive Auto-Switching", "Automatically routes only image-heavy pages to vision to optimize cost."),
        ("👁️ Token Grounding Layer", "Verifies AI structured fields against raw source tokens for zero hallucination.")
    ]
    for title, desc in s2_items:
        pt = tf_s2.add_paragraph()
        pt.space_before = Pt(8)
        pt.text = title
        pt.font.size = Pt(10.5)
        pt.font.bold = True
        pt.font.color.rgb = TEXT_BLACK

        pd = tf_s2.add_paragraph()
        pd.space_before = Pt(1)
        pd.text = f"    {desc}"
        pd.font.size = Pt(9.5)
        pd.font.color.rgb = TEXT_MUTED

    slide3.notes_slide.notes_text_frame.text = (
        "Here is the core technical differentiator. Traditional rule-based scrapers broke easily. "
        "Our architecture uses an adaptive dual-stream pipeline. When a document is a native digital PDF, we extract digital "
        "text at lightning speed. When encountering scanned documents, complex catalog diagrams, or handwritten logbooks, the system "
        "automatically triggers high-resolution image rendering and multimodal vision LLMs to parse the content with human-grade comprehension."
    )

    # ==========================================
    # SLIDE 4: CORE TECHNOLOGY STACK & SYSTEM FLOW
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide4, 4, 6)
    add_editorial_header(
        slide4,
        "03 // CORE TECHNOLOGY STACK & SYSTEM PIPELINE",
        "Core Technology Stack & System Execution Flow",
        "High-concurrency backend services, multimodal LLM orchestration, and high-fidelity parsing",
        tag_color=COLOR_VIOLET
    )

    # --- TOP: COLOR DIFFERENTIATED HORIZONTAL PIPELINE TIMELINE ---
    flow_y = Inches(2.0)
    flow_step_w = Inches(2.18)
    flow_step_gap = Inches(0.208)
    flow_x_start = Inches(0.8)

    flow_steps_4 = [
        ("01", "INTAKE", "PDF / SharePoint Intake", COLOR_VIOLET, "📄"),
        ("02", "ASYNC QUEUE", "FastAPI Worker Queue", COLOR_INDIGO, "⚡"),
        ("03", "PARSE & OCR", "PyMuPDF Text Stream", COLOR_TEAL, "⚙"),
        ("04", "AI EXTRACT", "Gemini 2.5 Vision Structuring", COLOR_EMERALD, "👁️"),
        ("05", "VALIDATE", "Grounding Sync & Export", COLOR_AMBER, "✓")
    ]

    for idx, (num, stage, desc, col, sym) in enumerate(flow_steps_4):
        step_x = flow_x_start + idx * (flow_step_w + flow_step_gap)

        add_hairline(slide4, step_x, flow_y, Inches(0.6), Inches(0.03), color=col)

        tb_step = slide4.shapes.add_textbox(step_x, flow_y + Inches(0.08), flow_step_w, Inches(0.85))
        tf_step = tb_step.text_frame
        tf_step.word_wrap = True
        tf_step.margin_left = tf_step.margin_top = tf_step.margin_right = tf_step.margin_bottom = 0

        p_b = tf_step.paragraphs[0]
        p_b.text = f"{sym} STAGE {num} //"
        p_b.font.size = Pt(8.5)
        p_b.font.bold = True
        p_b.font.color.rgb = col

        p_t = tf_step.add_paragraph()
        p_t.space_before = Pt(2)
        p_t.text = stage
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_BLACK

        p_d = tf_step.add_paragraph()
        p_d.space_before = Pt(1)
        p_d.text = desc
        p_d.font.size = Pt(9)
        p_d.font.color.rgb = TEXT_MUTED

        if idx < len(flow_steps_4) - 1:
            arr_x = step_x + flow_step_w + Inches(0.04)
            tb_arr = slide4.shapes.add_textbox(arr_x, flow_y + Inches(0.24), Inches(0.15), Inches(0.3))
            tf_arr = tb_arr.text_frame
            tf_arr.word_wrap = False
            tf_arr.margin_left = tf_arr.margin_top = tf_arr.margin_right = tf_arr.margin_bottom = 0
            pa = tf_arr.paragraphs[0]
            pa.text = "➔"
            pa.font.size = Pt(11)
            pa.font.bold = True
            pa.font.color.rgb = TEXT_MUTED

    # Horizontal Divider between Flow and Columns
    add_hairline(slide4, Inches(0.8), Inches(3.0), Inches(11.733), Inches(0.01))

    # --- BOTTOM: 4 COLOR-CODED TECHNOLOGY COLUMNS ---
    tech_y = Inches(3.18)
    tech_col_w = Inches(2.783)
    tech_col_gap = Inches(0.2)
    tech_h = Inches(3.5)

    tech_cards = [
        {
            "num": "01",
            "tag": "BACKEND API",
            "title": "Python FastAPI",
            "color": COLOR_VIOLET,
            "sym": "⚡",
            "bullets": [
                ("⚡ Async Architecture", "Non-blocking background job queue for high concurrency."),
                ("⚡ Polling Endpoints", "REST status endpoints with real-time extraction progress."),
                ("⚡ Type Validation", "Pydantic v2 schemas with auto OpenAPI documentation.")
            ]
        },
        {
            "num": "02",
            "tag": "MULTIMODAL AI",
            "title": "Gemini 2.5 Vision",
            "color": COLOR_TEAL,
            "sym": "👁️",
            "bullets": [
                ("👁️ Visual Reasoning", "Understands complex tabular hierarchies, notes, and callouts."),
                ("👁️ Handwriting / Notes", "Deciphers degraded scans, stamps, and field drawings."),
                ("👁️ Local Fallback", "Configurable Ollama local LLM execution for offline use.")
            ]
        },
        {
            "num": "03",
            "tag": "PARSING ENGINE",
            "title": "PyMuPDF (fitz)",
            "color": COLOR_EMERALD,
            "sym": "⚙",
            "bullets": [
                ("⚙ Sub-Second Parsing", "Direct digital glyph and embedded table extraction from PDFs."),
                ("⚙ 300-DPI Rasterizer", "High-resolution page rendering with auto-rotation."),
                ("⚙ Zero Latency Mode", "Bypasses OCR overhead completely when digital text exists.")
            ]
        },
        {
            "num": "04",
            "tag": "CLIENT & EXPORT",
            "title": "Vanilla JS & SheetJS",
            "color": COLOR_AMBER,
            "sym": "📊",
            "bullets": [
                ("📊 Zero Bundle Bloat", "Lightweight responsive UI designed for operational speed."),
                ("📊 Multi-Tab Excel", "Instant client-side generation of validated O&M workbooks."),
                ("📊 Cloud Integrations", "Direct webhooks for Microsoft Fabric, SharePoint, and ERP sync.")
            ]
        }
    ]

    for i, tc in enumerate(tech_cards):
        cx = flow_x_start + i * (tech_col_w + tech_col_gap)

        add_hairline(slide4, cx, tech_y, Inches(0.8), Inches(0.04), color=tc["color"])

        if i < len(tech_cards) - 1:
            add_hairline(slide4, cx + tech_col_w + Inches(0.1), tech_y, Inches(0.01), tech_h)

        tb_tc = slide4.shapes.add_textbox(cx, tech_y + Inches(0.1), tech_col_w, tech_h - Inches(0.1))
        tf_tc = tb_tc.text_frame
        tf_tc.word_wrap = True
        tf_tc.margin_left = tf_tc.margin_top = tf_tc.margin_right = tf_tc.margin_bottom = 0

        p0 = tf_tc.paragraphs[0]
        p0.text = f"{tc['sym']} {tc['num']} // {tc['tag']}"
        p0.font.size = Pt(8.5)
        p0.font.bold = True
        p0.font.color.rgb = tc["color"]

        p1 = tf_tc.add_paragraph()
        p1.space_before = Pt(3)
        p1.text = tc["title"]
        p1.font.size = Pt(13.5)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_BLACK

        for btitle, bdesc in tc["bullets"]:
            pt = tf_tc.add_paragraph()
            pt.space_before = Pt(8)
            pt.text = btitle
            pt.font.size = Pt(10)
            pt.font.bold = True
            pt.font.color.rgb = TEXT_BLACK

            pd = tf_tc.add_paragraph()
            pd.space_before = Pt(1)
            pd.text = f"    {bdesc}"
            pd.font.size = Pt(8.5)
            pd.font.color.rgb = TEXT_MUTED

    slide4.notes_slide.notes_text_frame.text = (
        "Slide 4 details our core technology stack and the end-to-end data execution pipeline. On the backend, we chose Python FastAPI for its native asynchronous capabilities and robust concurrency. For document manipulation, we use PyMuPDF for high-speed page parsing and rasterization. The extracted text or rendered image is sent to Gemini Multimodal LLM to generate structured schemas for maintenance routines and spare parts. This flow is managed through an asynchronous background queue, ensuring the system remains responsive even with large multi-hundred page manuals."
    )

    # ==========================================
    # SLIDE 5: AWS CLOUD ARCHITECTURE
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide5, 5, 6)
    add_editorial_header(
        slide5,
        "04 // CLOUD INFRASTRUCTURE & SCALABILITY",
        "AWS Cloud Architecture & Deployment Pipeline",
        "Enterprise-grade scalability, serverless container execution, and zero-trust security",
        tag_color=COLOR_INDIGO
    )

    # --- TOP: COLOR DIFFERENTIATED CLOUD TIMELINE ---
    flow_steps_5 = [
        ("01", "INGRESS", "HTTPS TLS 1.3 Client", COLOR_VIOLET, "🌐"),
        ("02", "TRAFFIC", "Route 53 & ALB WAF", COLOR_CYAN, "🔀"),
        ("03", "COMPUTE", "AWS ECS Fargate Fleet", COLOR_INDIGO, "⚙"),
        ("04", "STORAGE & AI", "Amazon S3 KMS ⮀ Gemini", COLOR_TEAL, "🗄️"),
        ("05", "DELIVERY", "Excel, DB & ERP Sync", COLOR_EMERALD, "🚀")
    ]

    for idx, (num, stage, desc, col, sym) in enumerate(flow_steps_5):
        step_x = flow_x_start + idx * (flow_step_w + flow_step_gap)

        add_hairline(slide5, step_x, flow_y, Inches(0.6), Inches(0.03), color=col)

        tb_step = slide5.shapes.add_textbox(step_x, flow_y + Inches(0.08), flow_step_w, Inches(0.85))
        tf_step = tb_step.text_frame
        tf_step.word_wrap = True
        tf_step.margin_left = tf_step.margin_top = tf_step.margin_right = tf_step.margin_bottom = 0

        p_b = tf_step.paragraphs[0]
        p_b.text = f"{sym} LAYER {num} //"
        p_b.font.size = Pt(8.5)
        p_b.font.bold = True
        p_b.font.color.rgb = col

        p_t = tf_step.add_paragraph()
        p_t.space_before = Pt(2)
        p_t.text = stage
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_BLACK

        p_d = tf_step.add_paragraph()
        p_d.space_before = Pt(1)
        p_d.text = desc
        p_d.font.size = Pt(9)
        p_d.font.color.rgb = TEXT_MUTED

        if idx < len(flow_steps_5) - 1:
            arr_x = step_x + flow_step_w + Inches(0.04)
            tb_arr = slide5.shapes.add_textbox(arr_x, flow_y + Inches(0.24), Inches(0.15), Inches(0.3))
            tf_arr = tb_arr.text_frame
            tf_arr.word_wrap = False
            tf_arr.margin_left = tf_arr.margin_top = tf_arr.margin_right = tf_arr.margin_bottom = 0
            pa = tf_arr.paragraphs[0]
            pa.text = "➔"
            pa.font.size = Pt(11)
            pa.font.bold = True
            pa.font.color.rgb = TEXT_MUTED

    # Horizontal Divider between Flow and Pillars
    add_hairline(slide5, Inches(0.8), Inches(3.0), Inches(11.733), Inches(0.01))

    # --- BOTTOM: 3 CLOUD ARCHITECTURE PILLARS ---
    cloud_col_w = Inches(3.644)
    cloud_col_gap = Inches(0.4)
    cloud_h = Inches(3.5)

    cloud_pillars = [
        {
            "num": "01",
            "tag": "SERVERLESS COMPUTE",
            "title": "AWS ECS & Fargate",
            "color": COLOR_INDIGO,
            "sym": "⚙",
            "bullets": [
                ("⚙ Auto-Scaling Worker Fleet", "Docker containers scale dynamically based on job queue volume."),
                ("⚙ Zero Server Overhead", "Serverless execution eliminates OS patching and capacity overhead."),
                ("⚙ Sandboxed Execution", "Each PDF extraction runs inside an isolated container instance.")
            ]
        },
        {
            "num": "02",
            "tag": "STORAGE & SECURITY",
            "title": "Amazon S3 & Secrets",
            "color": COLOR_TEAL,
            "sym": "🔒",
            "bullets": [
                ("🔒 Encrypted Buckets", "Raw PDFs, tiles, and JSON outputs encrypted at rest via SSE-KMS."),
                ("🔒 AWS Secrets Manager", "Centralized API key management with automated rotation for LLMs."),
                ("🔒 IAM Least-Privilege", "Strict execution policies ensuring zero-trust cloud security.")
            ]
        },
        {
            "num": "03",
            "tag": "ROUTING & MONITORING",
            "title": "ALB & CloudWatch",
            "color": COLOR_EMERALD,
            "sym": "📊",
            "bullets": [
                ("📊 Application Load Balancer", "SSL/TLS termination, health probes, and DDoS mitigation via AWS WAF."),
                ("📊 Real-Time Telemetry", "CloudWatch metrics for processing latency, error rates, and token usage."),
                ("📊 Automated Health Probes", "Proactive container restarts and instant alerts on failed tasks.")
            ]
        }
    ]

    for i, cp in enumerate(cloud_pillars):
        cx = flow_x_start + i * (cloud_col_w + cloud_col_gap)

        add_hairline(slide5, cx, tech_y, Inches(1.0), Inches(0.04), color=cp["color"])

        if i < len(cloud_pillars) - 1:
            add_hairline(slide5, cx + cloud_col_w + Inches(0.2), tech_y, Inches(0.01), cloud_h)

        tb_cp = slide5.shapes.add_textbox(cx, tech_y + Inches(0.1), cloud_col_w, cloud_h - Inches(0.1))
        tf_cp = tb_cp.text_frame
        tf_cp.word_wrap = True
        tf_cp.margin_left = tf_cp.margin_top = tf_cp.margin_right = tf_cp.margin_bottom = 0

        p0 = tf_cp.paragraphs[0]
        p0.text = f"{cp['sym']} {cp['num']} // {cp['tag']}"
        p0.font.size = Pt(9.0)
        p0.font.bold = True
        p0.font.color.rgb = cp["color"]

        p1 = tf_cp.add_paragraph()
        p1.space_before = Pt(3)
        p1.text = cp["title"]
        p1.font.size = Pt(14.5)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_BLACK

        for btitle, bdesc in cp["bullets"]:
            pt = tf_cp.add_paragraph()
            pt.space_before = Pt(8)
            pt.text = btitle
            pt.font.size = Pt(10.5)
            pt.font.bold = True
            pt.font.color.rgb = TEXT_BLACK

            pd = tf_cp.add_paragraph()
            pd.space_before = Pt(1)
            pd.text = f"    {bdesc}"
            pd.font.size = Pt(9)
            pd.font.color.rgb = TEXT_MUTED

    slide5.notes_slide.notes_text_frame.text = (
        "Slide 5 illustrates our AWS cloud deployment architecture. Traffic enters securely through Route 53 and an Application Load Balancer with SSL/TLS termination. The compute layer is powered by AWS ECS Fargate, running containerized FastAPI services in a serverless environment that auto-scales based on job queue depth. All documents and intermediate artifacts are stored securely in encrypted Amazon S3 buckets. Secret keys and IAM access are tightly controlled via AWS Secrets Manager, with full observability through CloudWatch metrics."
    )

    # ==========================================
    # SLIDE 6: THANK YOU & DISCUSSION SLIDE
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide6, 6, 6)

    # Micro Tag
    tb_ty_tag = slide6.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.733), Inches(0.3))
    tf_ty = tb_ty_tag.text_frame
    tf_ty.margin_left = tf_ty.margin_top = tf_ty.margin_right = tf_ty.margin_bottom = 0
    p_ty = tf_ty.paragraphs[0]
    p_ty.text = "✦  05 // CONCLUSION & NEXT STEPS"
    p_ty.font.size = Pt(10)
    p_ty.font.bold = True
    p_ty.font.color.rgb = COLOR_EMERALD

    # Giant Main Title
    tb_ty_title = slide6.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(11.733), Inches(1.8))
    tf_tt = tb_ty_title.text_frame
    tf_tt.word_wrap = True
    tf_tt.margin_left = tf_tt.margin_top = tf_tt.margin_right = tf_tt.margin_bottom = 0

    p_tt1 = tf_tt.paragraphs[0]
    p_tt1.text = "Thank You"
    p_tt1.font.size = Pt(44)
    p_tt1.font.bold = True
    p_tt1.font.color.rgb = TEXT_BLACK

    p_tt2 = tf_tt.add_paragraph()
    p_tt2.space_before = Pt(4)
    p_tt2.text = "Questions & Technical Discussion"
    p_tt2.font.size = Pt(18)
    p_tt2.font.bold = True
    p_tt2.font.color.rgb = COLOR_TEAL

    # Architectural Mid Divider Line
    add_hairline(slide6, Inches(0.8), Inches(3.65), Inches(11.733), Inches(0.015), color=LINE_DARK)

    # Two Column Layout below title
    # Left Column: Platform Summary
    tb_ty_left = slide6.shapes.add_textbox(Inches(0.8), Inches(3.95), Inches(5.6), Inches(2.3))
    tf_ty_left = tb_ty_left.text_frame
    tf_ty_left.word_wrap = True
    tf_ty_left.margin_left = tf_ty_left.margin_top = tf_ty_left.margin_right = tf_ty_left.margin_bottom = 0

    p_tl1 = tf_ty_left.paragraphs[0]
    p_tl1.text = "Intelligent Document Processing (IDP) Platform"
    p_tl1.font.size = Pt(15)
    p_tl1.font.bold = True
    p_tl1.font.color.rgb = TEXT_HEAVY

    p_tl2 = tf_ty_left.add_paragraph()
    p_tl2.space_before = Pt(10)
    p_tl2.text = "Enterprise-grade pipeline designed for high-concurrency ingestion of complex O&M manuals, vendor catalogs, and handwritten shift logs into structured, verified asset data."
    p_tl2.font.size = Pt(11.5)
    p_tl2.font.color.rgb = TEXT_MUTED

    # Vertical Hairline Divider
    add_hairline(slide6, Inches(6.8), Inches(3.95), Inches(0.01), Inches(2.4))

    # Right Column: Discussion Points
    tb_ty_right = slide6.shapes.add_textbox(Inches(7.2), Inches(3.95), Inches(5.3), Inches(2.4))
    tf_ty_right = tb_ty_right.text_frame
    tf_ty_right.word_wrap = True
    tf_ty_right.margin_left = tf_ty_right.margin_top = tf_ty_right.margin_right = tf_ty_right.margin_bottom = 0

    ty_caps = [
        ("✦ DISCUSSION & COLLABORATION", "Open for architectural questions, pilot integration, and technical deep-dives.", COLOR_VIOLET),
        ("⚙ INTEGRATION WORKFLOW", "Seamless connectivity with Microsoft SharePoint, Fabric Lakehouse, and ERP asset registries.", COLOR_TEAL),
        ("🚀 NEXT STEPS", "Live sandbox demonstration, custom schema onboarding, and throughput benchmarking.", COLOR_EMERALD)
    ]

    for idx, (ty_title, ty_desc, ty_col) in enumerate(ty_caps):
        p_c_title = tf_ty_right.paragraphs[0] if idx == 0 else tf_ty_right.add_paragraph()
        if idx > 0:
            p_c_title.space_before = Pt(10)
        p_c_title.text = ty_title
        p_c_title.font.size = Pt(10.5)
        p_c_title.font.bold = True
        p_c_title.font.color.rgb = ty_col

        p_c_desc = tf_ty_right.add_paragraph()
        p_c_desc.space_before = Pt(2)
        p_c_desc.text = ty_desc
        p_c_desc.font.size = Pt(10)
        p_c_desc.font.color.rgb = TEXT_MUTED

    slide6.notes_slide.notes_text_frame.text = (
        "Thank you all for your time and attention today. We have covered the use case journey, our dual-stream parsing engine, "
        "the FastAPI and Gemini technology stack, and our scalable AWS cloud deployment. I would now love to open the floor for any questions, "
        "technical discussions, or pilot integration opportunities."
    )

    output_path = "/Users/akshayryali/1/IDP_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
